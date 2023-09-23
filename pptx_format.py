#!/usr/bin/env python
from pptx import Presentation
from tools import RetVal
from tools import Settings
import os

class Pptx_Format():

    def init(self, in_path="", out_path= "") -> None:
         self.input_dir = in_path
         self.output_dir = out_path

    def CountFiles(self):
         dir_list = os.listdir(self.input_dir)
         dir_list_pptx = [file for file in dir_list if file.endswith(".pptx")]
         return len(dir_list_pptx)
    
    def ListFiles(self):
         dir_list = os.listdir(self.input_dir)
         dir_list_pptx = [file for file in dir_list if file.endswith(".pptx")]
         return dir_list_pptx
    
    def ProcessFiles(self):
        retVal = RetVal.ERROR
         
        print("Processing pptx files...")
        if self.CountFiles() == 0:
            print("ERROR: Skipping PPTX processing, no files found!")
            return retVal
        
        pptx_file_list = self.ListFiles()
        
        for i in range(0, len(pptx_file_list)):
        # for i in range(0, 1):
            print(f"    > File {i + 1}/{len(pptx_file_list)}", end='')

            #Extract song title
            song_title = Settings.SONG_PREFIX + pptx_file_list[i].split('.')[0]

            #Open input file
            input_file = open(self.input_dir + '/' + pptx_file_list[i], "rb")

            #Open output file
            output_file = open(f"{self.output_dir}/{song_title}.txt", "w+",encoding="utf-8")

            #Open powrpoint presentation
            prs = Presentation(input_file)

            #Print the song title tag
            output_file.write(f"{{title: {song_title}}}\n")

            #Add the author Tag
            if Settings.AUTHOR != "":
                output_file.write(f"{{author: {Settings.AUTHOR}}}\n")

            #Parse each slide (each slide represents 1 verse)
            verse = 1

            for slide in prs.slides:
                output_file.write(f"\n{{comment: Verse {verse}}}\n")
                
                #Parse each shape in slide and filter for text ones
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue

                    #For each paragraph found in shape
                    for paragraph in shape.text_frame.paragraphs:
                        #For each sentence in paragraph
                        for sentence in paragraph.runs:
                            #Filter for small words, like slide 1/3
                            if len(sentence.text) > Settings.MIN_SENTENCE_LEN and sentence.text not in Settings.ignored_keywords:
                                output_file.write(sentence.text + '\n')

                #Increment verse for next slide
                verse += 1

            #Close input and output files
            print(" - DONE")
            input_file.close()
            output_file.close()

        print("DONE processing pptx files")
        retVal = RetVal.OK

        return retVal